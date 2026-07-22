"""Backend adapters — the thin shims that make a heterogeneous backend speak the
one contract. Each returns raw `(outputs, runtime, status, error, degraded)`;
the gateway seals the receipt and types the warrant. New compute kind = new
adapter here + a registry row.

Adapters are injectable (`set_backend`) so tests never need a live forge/graph.
"""
from __future__ import annotations

import base64
import hashlib
import io
import os
import re
import sqlite3
import time
from typing import Any, Awaitable, Callable

import httpx

from .contract import ComputeOutput, GraphDelta, GraphEdge, GraphNode

FORGE_URL = os.getenv("FORGE_URL", "http://lattice-forge.sovereign-runtime.svc.cluster.local:8870").rstrip("/")
FORGE_TOKEN = os.getenv("FORGE_TOKEN", "")
HELLGRAPH_URL = os.getenv("HELLGRAPH_URL", "http://hellgraph-service:8090").rstrip("/")
SPARK_RUNNER_URL = os.getenv("SPARK_RUNNER_URL", "http://spark-runner:8080").rstrip("/")
SPARK_RUNNER_TOKEN = os.getenv("SPARK_RUNNER_TOKEN", "")
MODEL_SERVER_URL = os.getenv("MODEL_SERVER_URL", "http://embeddings:8080").rstrip("/")
MODEL_SERVER_TOKEN = os.getenv("MODEL_SERVER_TOKEN", "")
# document extraction (Holmes/Sherlock via the Studio BFF) + the open-data reference source.
EXTRACT_URL = os.getenv("EXTRACT_URL", "http://lattice-studio:8080").rstrip("/")
EXTRACT_TOKEN = os.getenv("EXTRACT_TOKEN", "")
# SEC EDGAR XBRL company-facts — public, keyless; the FactSet stand-in for reconciliation.
SEC_EDGAR_URL = os.getenv("SEC_EDGAR_URL", "https://data.sec.gov").rstrip("/")
SEC_EDGAR_UA = os.getenv("SEC_EDGAR_UA", "SocioProphet compute-gateway compliance@socioprophet.dev")
TIMEOUT = float(os.getenv("GATEWAY_TIMEOUT", "90"))

# adapter result shape: dict(outputs=[ComputeOutput], runtime, status, error, degraded)
AdapterResult = dict[str, Any]


async def _forge(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    body = {"project": project, "code": req_spec.get("code", ""),
            "language": req_spec.get("language", "python"),
            "adapter": req_spec.get("adapter"), "session_id": session}
    headers = {"Authorization": f"Bearer {FORGE_TOKEN}"} if FORGE_TOKEN else {}
    try:
        async with httpx.AsyncClient(timeout=TIMEOUT, headers=headers) as c:
            r = await c.post(f"{FORGE_URL}/v1/execute", json=body)
            if r.status_code != 200:
                return {"outputs": [], "runtime": "python3", "status": "error",
                        "error": f"forge HTTP {r.status_code}", "degraded": None}
            d = r.json()
            outs = [ComputeOutput(**o) if not isinstance(o, ComputeOutput) else o
                    for o in _shape_forge(d.get("outputs", []))]
            return {"outputs": outs, "runtime": d.get("runtime", "python3"),
                    "status": d.get("status", "ok"), "error": d.get("error"),
                    "degraded": d.get("degraded")}
    except Exception as e:  # noqa: BLE001
        return {"outputs": [], "runtime": "python3", "status": "degraded",
                "error": None, "degraded": f"forge unreachable: {e}"}


def _shape_forge(outputs: list[dict]) -> list[dict]:
    shaped = []
    for o in outputs:
        shaped.append({"type": o.get("type", "result"), "text": o.get("text"),
                       "data": {k: o[k] for k in ("png", "svg", "html", "mime") if o.get(k)} or None,
                       "mime": o.get("mime")})
    return shaped


async def _hellgraph_query(req_spec: dict, project: str) -> AdapterResult:
    label = req_spec.get("label") or project
    try:
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.get(f"{HELLGRAPH_URL}/api/graph/query", params={"label": label})
            if r.status_code != 200:
                return {"outputs": [], "runtime": "hellgraph", "status": "error",
                        "error": f"hellgraph HTTP {r.status_code}", "degraded": None}
            data = r.json()
            nodes = data.get("nodes", data if isinstance(data, list) else [])
            return {"outputs": [ComputeOutput(type="graph", data={"nodes": nodes, "count": len(nodes)})],
                    "runtime": "hellgraph", "status": "ok", "error": None, "degraded": None}
    except Exception as e:  # noqa: BLE001
        return {"outputs": [], "runtime": "hellgraph", "status": "degraded",
                "error": None, "degraded": f"hellgraph unreachable: {e}"}


async def _hellgraph_stats(req_spec: dict, project: str) -> AdapterResult:
    try:
        async with httpx.AsyncClient(timeout=TIMEOUT) as c:
            r = await c.get(f"{HELLGRAPH_URL}/api/graph/stats")
            if r.status_code != 200:
                return {"outputs": [], "runtime": "hellgraph", "status": "error",
                        "error": f"hellgraph HTTP {r.status_code}", "degraded": None}
            return {"outputs": [ComputeOutput(type="table", data=r.json())],
                    "runtime": "hellgraph", "status": "ok", "error": None, "degraded": None}
    except Exception as e:  # noqa: BLE001
        return {"outputs": [], "runtime": "hellgraph", "status": "degraded",
                "error": None, "degraded": f"hellgraph unreachable: {e}"}


async def _spark(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Sovereign Spark: submit SQL/DataFrame code to spark-runner (the same runtime
    lattice-studio dispatches to). Databricks' one paradigm, here as one backend
    among many behind the uniform contract — entitlement-gated, receipt-sealed."""
    body = {"sql": req_spec.get("sql", ""), "data": req_spec.get("data", []),
            "table": req_spec.get("table", "t"), "job_id": session}
    headers = {"Authorization": f"Bearer {SPARK_RUNNER_TOKEN}"} if SPARK_RUNNER_TOKEN else {}
    try:
        async with httpx.AsyncClient(timeout=TIMEOUT, headers=headers) as c:
            r = await c.post(f"{SPARK_RUNNER_URL}/v1/submit", json=body)
            if r.status_code != 200:
                return {"outputs": [], "runtime": "spark", "status": "error",
                        "error": f"spark-runner HTTP {r.status_code}: {r.text[:200]}", "degraded": None}
            d = r.json()
            return {"outputs": [ComputeOutput(type="table", data={
                        "rows": d.get("rows", []), "row_count": d.get("row_count"),
                        "backend_receipt": d.get("receipt")})],
                    "runtime": "spark", "status": "ok", "error": None, "degraded": None}
    except Exception as e:  # noqa: BLE001
        return {"outputs": [], "runtime": "spark", "status": "degraded",
                "error": None, "degraded": f"spark-runner unreachable: {e}"}


async def _inference(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Model inference (embed | chat) via the sovereign model server. Output warrant
    is `derived` — a model produces it, it is not observed from the graph."""
    task = req_spec.get("task", "embed")
    if task == "embed":
        payload = {"input": req_spec.get("input") or req_spec.get("texts") or []}
        path = "/embed"
    else:
        payload = {"messages": req_spec.get("messages", []), "model": req_spec.get("model")}
        path = "/v1/chat/completions"
    headers = {"Authorization": f"Bearer {MODEL_SERVER_TOKEN}"} if MODEL_SERVER_TOKEN else {}
    try:
        async with httpx.AsyncClient(timeout=TIMEOUT, headers=headers) as c:
            r = await c.post(f"{MODEL_SERVER_URL}{path}", json=payload)
            if r.status_code != 200:
                return {"outputs": [], "runtime": "model-server", "status": "error",
                        "error": f"model-server HTTP {r.status_code}", "degraded": None}
            return {"outputs": [ComputeOutput(type="result", data=r.json())],
                    "runtime": "model-server", "status": "ok", "error": None, "degraded": None}
    except Exception as e:  # noqa: BLE001
        return {"outputs": [], "runtime": "model-server", "status": "degraded",
                "error": None, "degraded": f"model-server unreachable: {e}"}


# ── IFM stages 01–02: land the pack, parse it to blocks ──
_EMU_PER_PT = 12700  # PowerPoint positions are EMU; points are what humans cite


def _sniff_media(filename: str, raw: bytes) -> str:
    """Magic-bytes first (filenames lie), extension as tiebreak for zip containers."""
    if raw[:5] == b"%PDF-":
        return "application/pdf"
    if raw[:4] == b"PK\x03\x04":  # OOXML container — pptx and docx share it
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext == "pptx":
            return "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if ext == "docx":
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        return "application/zip"
    return "text/plain"


async def _ingest(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Stage 01 — land the pack by content hash. The sha256 is over the RAW bytes, so
    identical packs hash identically whatever they're named, and a re-run of the same
    request memoizes (the dedupe). Threads document_b64 + media_type to parse."""
    try:
        raw = base64.b64decode(req_spec.get("document_b64") or "", validate=True)
    except Exception:  # noqa: BLE001
        return {"outputs": [], "runtime": "gateway", "status": "error",
                "error": "ingest needs document_b64 (valid base64)", "degraded": None}
    if not raw:
        return {"outputs": [], "runtime": "gateway", "status": "error",
                "error": "ingest needs document_b64 (valid base64)", "degraded": None}
    filename = str(req_spec.get("filename") or "")
    media = str(req_spec.get("media_type") or _sniff_media(filename, raw))
    return {"outputs": [ComputeOutput(type="artifact", data={
                "sha256": hashlib.sha256(raw).hexdigest(), "size": len(raw),
                "media_type": media, "filename": filename,
                "document_b64": req_spec.get("document_b64")})],
            "runtime": "gateway", "status": "ok", "error": None, "degraded": None,
            "epistemic": "observed"}


def _pdf_blocks(raw: bytes) -> tuple[list[dict], int]:
    from pypdf import PdfReader  # lazy: parse is the only kind that needs it
    reader = PdfReader(io.BytesIO(raw))
    blocks: list[dict] = []
    for pno, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        region = 0
        for para in re.split(r"\n\s*\n", text):
            para = para.strip()
            if para:
                blocks.append({"page": pno, "region": region, "kind": "text", "text": para})
                region += 1
    return blocks, len(reader.pages)


def _pptx_blocks(raw: bytes) -> tuple[list[dict], int]:
    from pptx import Presentation  # lazy, same reason
    prs = Presentation(io.BytesIO(raw))
    blocks: list[dict] = []
    n = 0
    for sno, slide in enumerate(prs.slides, 1):
        n = sno
        for shape in slide.shapes:
            bbox = None
            if shape.left is not None and shape.top is not None:
                bbox = [round(shape.left / _EMU_PER_PT, 1), round(shape.top / _EMU_PER_PT, 1),
                        round((shape.width or 0) / _EMU_PER_PT, 1), round((shape.height or 0) / _EMU_PER_PT, 1)]
            if getattr(shape, "has_table", False):
                # tables carry the numbers — cells pipe-joined row-wise, one block per table
                rows = "\n".join(" | ".join(c.text.strip() for c in row.cells) for row in shape.table.rows)
                if rows.strip():
                    blocks.append({"page": sno, "kind": "table", "text": rows, "bbox": bbox})
            elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                blocks.append({"page": sno, "kind": "text", "text": shape.text_frame.text.strip(), "bbox": bbox})
    return blocks, n


async def _parse(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Stage 02 — document bytes → blocks[], each keeping its page (+ bbox in points for
    PPTX shapes, tables pipe-joined; page/paragraph regions for PDF). Consumes the fields
    ingest threads in; a corrupt or unsupported document is an ERROR, never a crash."""
    try:
        raw = base64.b64decode(req_spec.get("document_b64") or "", validate=True)
    except Exception:  # noqa: BLE001
        return {"outputs": [], "runtime": "gateway", "status": "error",
                "error": "parse needs document_b64 (thread it from an ingest step)", "degraded": None}
    if not raw:
        return {"outputs": [], "runtime": "gateway", "status": "error",
                "error": "parse needs document_b64 (thread it from an ingest step)", "degraded": None}
    filename = str(req_spec.get("filename") or "")
    media = str(req_spec.get("media_type") or _sniff_media(filename, raw))
    try:
        if media == "application/pdf":
            blocks, pages = _pdf_blocks(raw)
        elif media.endswith("presentationml.presentation"):
            blocks, pages = _pptx_blocks(raw)
        elif media.startswith("text/"):
            text = raw.decode("utf-8", errors="replace").strip()
            blocks = [{"page": 1, "region": i, "kind": "text", "text": p.strip()}
                      for i, p in enumerate(re.split(r"\n\s*\n", text)) if p.strip()]
            pages = 1
        else:
            return {"outputs": [], "runtime": "gateway", "status": "error",
                    "error": f"parse: unsupported media type {media}", "degraded": None}
    except Exception as e:  # noqa: BLE001 — corrupt document → honest error, not a 500
        return {"outputs": [], "runtime": "gateway", "status": "error",
                "error": f"parse failed: {e}", "degraded": None}
    if not blocks:
        return {"outputs": [], "runtime": "gateway", "status": "error",
                "error": "parse: no extractable content (image-only/scanned documents need OCR)",
                "degraded": None}
    return {"outputs": [ComputeOutput(type="blocks", data={
                "blocks": blocks, "pages": pages, "media_type": media,
                "sha256": req_spec.get("sha256"), "filename": filename})],
            "runtime": "gateway", "status": "ok", "error": None, "degraded": None,
            "epistemic": "observed"}


# ── IFM: document → structured facts (extraction) + reconcile vs open data ──
# Each extracted fact carries its lineage AND an epistemic warrant, so the desk can
# gate what it trades on. observed = verbatim on the page; derived = computed.
_WARRANT_ORDER = ["unknown", "hypothesis", "simulated", "observed", "derived", "verified", "attested"]


def _norm_fact(f: dict) -> dict:
    """Normalise an extracted fact into the ExtractedFact shape + default warrant."""
    warrant = f.get("warrant")
    if warrant not in _WARRANT_ORDER:
        # verbatim-with-a-source-span is `observed`; anything computed is `derived`
        warrant = "observed" if (f.get("verbatim") or f.get("source_span")) else "derived"
    return {"field": f.get("field"), "value": f.get("value"), "unit": f.get("unit"),
            "page": f.get("page"), "source_span": f.get("source_span"),
            "confidence": f.get("confidence"), "warrant": warrant}


async def _extraction(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Extract a document into typed rows against a target schema (Holmes/Sherlock via
    the Studio extract path). Output warrant = the WEAKEST fact's warrant. If the spec
    already carries `facts` (pre-parsed), extract is a no-op typing pass — useful for a
    demo pack whose blocks were parsed upstream."""
    schema = req_spec.get("target_schema", {})
    facts = req_spec.get("facts")
    if facts is None:
        headers = {"Authorization": f"Bearer {EXTRACT_TOKEN}"} if EXTRACT_TOKEN else {}
        body = {"project": project, "document": req_spec.get("document"),
                "blocks": req_spec.get("blocks"), "target_schema": schema,
                "period": req_spec.get("period")}
        try:
            async with httpx.AsyncClient(timeout=TIMEOUT, headers=headers) as c:
                # fact-mode endpoint: {blocks, target_schema, period} → {facts[]} — NOT the
                # graph-writing /api/studio/extract, whose contract is text→entities.
                r = await c.post(f"{EXTRACT_URL}/api/studio/extract-facts", json=body)
                if r.status_code != 200:
                    return {"outputs": [], "runtime": "holmes", "status": "error",
                            "error": f"extract HTTP {r.status_code}", "degraded": None}
                facts = r.json().get("facts", [])
        except Exception as e:  # noqa: BLE001
            return {"outputs": [], "runtime": "holmes", "status": "degraded", "error": None,
                    "degraded": f"extractor unreachable: {e}"}
    rows = [_norm_fact(f) for f in facts]
    warrant = min((r["warrant"] for r in rows),
                  key=lambda w: _WARRANT_ORDER.index(w)) if rows else "unknown"
    return {"outputs": [ComputeOutput(type="table", data={
                "table": schema.get("table"), "rows": rows,
                "entity": req_spec.get("entity"), "period": req_spec.get("period")})],
            "runtime": "holmes", "status": "ok", "error": None, "degraded": None,
            "epistemic": warrant}


# field → us-gaap XBRL concept (extend as the schema grows)
_XBRL_CONCEPT = {
    "revenue": "RevenueFromContractWithCustomerExcludingAssessedTax",
    "revenues": "Revenues", "net_income": "NetIncomeLoss", "net_profit": "NetIncomeLoss",
    "npat": "NetIncomeLoss",
    "gross_profit": "GrossProfit", "cost_of_revenue": "CostOfRevenue",
    "operating_income": "OperatingIncomeLoss", "operating_profit": "OperatingIncomeLoss",
    "eps_basic": "EarningsPerShareBasic", "eps_diluted": "EarningsPerShareDiluted",
    "assets": "Assets", "liabilities": "Liabilities", "equity": "StockholdersEquity",
    "cash": "CashAndCashEquivalentsAtCarryingValue",
    "operating_cash_flow": "NetCashProvidedByUsedInOperatingActivities",
}

# companyfacts returns EVERY concept for an entity in one (large) payload — cached per
# CIK so a multi-field reconcile is one EDGAR call, not one per field.
_EDGAR_CACHE: dict[str, tuple[float, dict]] = {}
_EDGAR_CACHE_TTL = float(os.getenv("SEC_EDGAR_CACHE_TTL", "900"))


async def _edgar_companyfacts(cik: str) -> dict | None:
    hit = _EDGAR_CACHE.get(cik)
    if hit and (time.time() - hit[0]) < _EDGAR_CACHE_TTL:
        return hit[1]
    try:
        async with httpx.AsyncClient(timeout=TIMEOUT, headers={"User-Agent": SEC_EDGAR_UA}) as c:
            r = await c.get(f"{SEC_EDGAR_URL}/api/xbrl/companyfacts/CIK{cik}.json")
            if r.status_code != 200:
                return None
            data = r.json()
    except Exception:  # noqa: BLE001
        return None
    _EDGAR_CACHE[cik] = (time.time(), data)
    return data


async def _sec_edgar_reference(entity: dict, field: str, period: str) -> float | None:
    """US reference — SEC EDGAR XBRL via the company-facts endpoint the design names
    (public, keyless; one call per entity, cached). `entity` carries a zero-padded 10-digit
    `cik`. Returns the value whose fiscal period matches, else None (unresolved → the fact
    simply can't reach `verified`)."""
    cik = str(entity.get("cik", "")).zfill(10)
    concept = _XBRL_CONCEPT.get(field.lower())
    if not cik.strip("0") or concept is None:
        return None
    data = await _edgar_companyfacts(cik)
    fact = ((data or {}).get("facts", {}).get("us-gaap", {}) or {}).get(concept) or {}
    for unit_rows in (fact.get("units") or {}).values():
        for row in unit_rows:
            if period and (row.get("fp") == period or str(row.get("frame", "")).find(period) >= 0
                           or str(row.get("fy")) == period):
                try:
                    return float(row["val"])
                except (KeyError, TypeError, ValueError):
                    continue
    return None


async def _asx_prior_extraction_reference(entity: dict, field: str, period: str) -> float | None:
    """AU reference — cross-document: Australia has no open-XBRL EDGAR twin (ASIC is paid),
    but every ASX results release is lodged WITH a statutory Appendix 4E/4D. Run the
    statutory form through this same pipeline FIRST (its rows land in the SQL sink,
    receipts and all), then the glossy investor pack reconciles against those rows —
    reference data that itself carries provenance. Reads (entity, period, field) from the
    reference table in the load sink."""
    ent_id = str(entity.get("cik") or entity.get("asx") or entity.get("name") or "")
    table = str(entity.get("reference_table") or "reference_facts")
    if not ent_id or not table.replace("_", "").isalnum():
        return None
    try:
        con = sqlite3.connect(_sqlite_path(SQL_LOAD_DSN))
        row = con.execute(f"SELECT value_abs, value FROM {table} WHERE entity=? AND period=? AND field=?",
                          (ent_id, period, field)).fetchone()
        con.close()
    except Exception:  # noqa: BLE001
        return None
    if row is None:
        return None
    # absolute value preferred — reconcile compares in absolute units on both sides
    for v in row:
        try:
            if v is not None:
                return float(v)
        except (TypeError, ValueError):
            continue
    return None


_REFERENCE_BACKENDS: dict[str, Callable[..., Awaitable[float | None]]] = {
    "sec-edgar": _sec_edgar_reference,
    "asx-appendix": _asx_prior_extraction_reference,
}


def _pick_reference(source: str | None, entity: dict) -> tuple[str, Callable[..., Awaitable[float | None]]]:
    """Jurisdiction routing: an explicit source wins; else a CIK routes US (EDGAR) and an
    ASX ticker routes AU (statutory cross-document). FactSet lands here as one more entry —
    the reconcile logic never changes."""
    if source in _REFERENCE_BACKENDS:
        return source, _REFERENCE_BACKENDS[source]
    if entity.get("asx") and not entity.get("cik"):
        return "asx-appendix", _REFERENCE_BACKENDS["asx-appendix"]
    return "sec-edgar", _REFERENCE_BACKENDS["sec-edgar"]


# test override: when set, it beats jurisdiction routing (tests supply a deterministic source)
_REFERENCE: Callable[..., Awaitable[float | None]] | None = None


def set_reference_resolver(fn) -> None:
    global _REFERENCE
    _REFERENCE = fn


# magnitude suffixes a pack prints beside a number ('$1,204m', '2.5bn'). A unit string
# may carry a currency prefix ('AUD_m') — the suffix after the final '_' is what scales.
_UNIT_SCALE = {"k": 1e3, "m": 1e6, "b": 1e9, "bn": 1e9}


def _fact_abs_value(val: Any, unit: Any) -> float | None:
    """Printed value → absolute units for comparison (1204 + 'm' → 1_204_000_000).
    References (EDGAR XBRL, FactSet) speak absolute; packs print magnitudes — comparing
    them raw would false-flag every correct fact. Unknown/absent unit → as printed."""
    if val is None:
        return None
    try:
        v = float(val)
    except (TypeError, ValueError):
        return None
    u = str(unit or "").lower().rsplit("_", 1)[-1]
    return v * _UNIT_SCALE.get(u, 1.0)


async def _reconcile(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Reconcile extracted facts against a structured reference source (SEC EDGAR open
    data now; FactSet on a key later — the logic is identical). Agreement within tolerance
    PROMOTES a fact to `verified`; divergence is FLAGGED (the tradable edge). The step's
    warrant is `verified` iff every fact reconciled, else `derived`. Comparison happens in
    ABSOLUTE units (the design's validate-stage unit check): '$1,204m' vs XBRL 1204000000
    agree; a magnitude slip false-flags instead of silently verifying."""
    # `rows` is what the extraction step emits — accept it so a workflow can thread
    # extraction → reconcile via `from` without a rename.
    facts = req_spec.get("facts") or req_spec.get("rows") or []
    entity = req_spec.get("entity", {})
    period = req_spec.get("period", "")
    tol = float(req_spec.get("tolerance", 0.01))   # 1% default
    source, resolver = _pick_reference(req_spec.get("source"), entity if isinstance(entity, dict) else {})
    if _REFERENCE is not None:
        resolver = _REFERENCE               # injected test resolver beats routing
    recs, all_ok = [], True
    for f in facts:
        field, val, unit = f.get("field"), f.get("value"), f.get("unit")
        ref = await resolver(entity, field, period)
        val_abs = _fact_abs_value(val, unit)
        within = ref is not None and val_abs is not None and abs(val_abs - ref) <= abs(ref) * tol
        delta = (val_abs - ref) if (ref is not None and val_abs is not None) else None
        warrant = "verified" if within else (f.get("warrant") or "derived")
        all_ok = all_ok and within
        recs.append({"field": field, "extracted": val, "extracted_abs": val_abs, "unit": unit,
                     "reference": ref, "delta": delta,
                     "within_tol": within, "warrant": warrant, "flagged": ref is not None and not within})
    return {"outputs": [ComputeOutput(type="table", data={
                "entity": entity, "period": period, "source": source,
                "reconciliations": recs, "all_verified": all_ok and bool(recs)})],
            "runtime": source, "status": "ok", "error": None, "degraded": None,
            "epistemic": "verified" if (all_ok and recs) else "derived"}


# ── IFM: load reconciled facts into the structured SQL layer (the doc→SQL sink) ──
SQL_LOAD_DSN = os.getenv("SQL_LOAD_DSN", "sqlite:////tmp/ifm_extract.db")


def _sqlite_path(dsn: str) -> str:
    return dsn[len("sqlite:///"):] if dsn.startswith("sqlite:///") else dsn


async def _sql_load(req_spec: dict, project: str, session: str | None) -> AdapterResult:
    """Upsert reconciled facts into the consolidated SQL layer, keyed by (entity, period,
    field), each row carrying its warrant + reference + source. Demonstrated against a
    sovereign SQLite file (no creds); a Postgres DSN is the production swap. The load
    step's warrant is the WEAKEST row loaded — you can't trust the table more than its
    least-trustworthy cell."""
    rows = req_spec.get("rows") or req_spec.get("reconciliations") or req_spec.get("facts") or []
    table = req_spec.get("table") or (req_spec.get("target_schema") or {}).get("table") or "extracted_facts"
    if not table.replace("_", "").isalnum():
        return {"outputs": [], "runtime": "sql", "status": "error",
                "error": f"unsafe table name: {table!r}", "degraded": None}
    entity = req_spec.get("entity", {})
    # keying must match what the reference resolvers use: cik (US) → asx (AU) → name
    ent_id = str(entity.get("cik") or entity.get("asx") or entity.get("name") or entity) if isinstance(entity, dict) else str(entity)
    period, source = str(req_spec.get("period", "")), str(req_spec.get("source", "extraction"))
    if not rows:
        return {"outputs": [], "runtime": "sql", "status": "degraded", "error": None,
                "degraded": "no rows to load"}
    dsn = req_spec.get("dsn") or SQL_LOAD_DSN
    if not dsn.startswith("sqlite"):
        # Postgres/other = the production sink; the driver isn't vendored in the gateway.
        return {"outputs": [], "runtime": "sql", "status": "degraded", "error": None,
                "degraded": f"non-sqlite DSN needs the production driver (swap here): {dsn.split('://')[0]}"}
    warrants = []
    try:
        con = sqlite3.connect(_sqlite_path(dsn))
        con.execute(f"CREATE TABLE IF NOT EXISTS {table} (entity TEXT, period TEXT, field TEXT, "
                    "value TEXT, warrant TEXT, reference TEXT, within_tol INTEGER, source TEXT, "
                    "loaded_at TEXT, value_abs REAL, unit TEXT, PRIMARY KEY(entity, period, field))")
        for col, typ in (("value_abs", "REAL"), ("unit", "TEXT")):
            try:  # pre-normalization tables lack these columns — migrate in place
                con.execute(f"ALTER TABLE {table} ADD COLUMN {col} {typ}")
            except sqlite3.OperationalError:
                pass
        inserted = updated = 0
        now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
        for r in rows:
            field = r.get("field")
            value = r.get("extracted", r.get("value"))
            unit = r.get("unit")
            # absolute value beside the printed one, so cross-document reconciliation
            # (the AU path reads this table as its reference) compares like with like
            value_abs = r.get("extracted_abs")
            if value_abs is None:
                value_abs = _fact_abs_value(value, unit)
            warrant = r.get("warrant") or "derived"
            warrants.append(warrant)
            exists = con.execute(f"SELECT 1 FROM {table} WHERE entity=? AND period=? AND field=?",
                                 (ent_id, period, field)).fetchone() is not None
            con.execute(
                f"INSERT INTO {table}(entity,period,field,value,warrant,reference,within_tol,source,loaded_at,value_abs,unit) "
                "VALUES(?,?,?,?,?,?,?,?,?,?,?) ON CONFLICT(entity,period,field) DO UPDATE SET "
                "value=excluded.value, warrant=excluded.warrant, reference=excluded.reference, "
                "within_tol=excluded.within_tol, source=excluded.source, loaded_at=excluded.loaded_at, "
                "value_abs=excluded.value_abs, unit=excluded.unit",
                (ent_id, period, field, str(value), warrant,
                 str(r.get("reference")), 1 if r.get("within_tol") else 0, source, now,
                 value_abs, str(unit) if unit is not None else None))
            inserted, updated = (inserted, updated + 1) if exists else (inserted + 1, updated)
        con.commit()
        total = con.execute(f"SELECT COUNT(*) FROM {table}").fetchone()[0]
        con.close()
    except Exception as e:  # noqa: BLE001
        return {"outputs": [], "runtime": "sql", "status": "error",
                "error": f"sql load failed: {e}", "degraded": None}
    weakest = min(warrants, key=lambda w: _WARRANT_ORDER.index(w) if w in _WARRANT_ORDER else 0)
    return {"outputs": [ComputeOutput(type="table", data={
                "table": table, "target": dsn, "inserted": inserted, "updated": updated,
                "rows_written": inserted + updated, "table_total": total,
                "entity": ent_id, "period": period})],
            "runtime": "sql", "status": "ok", "error": None, "degraded": None, "epistemic": weakest}


# kind → adapter coroutine. Overridable in tests.
_BACKENDS: dict[str, Callable[..., Awaitable[AdapterResult]]] = {
    "forge": lambda spec, project, session: _forge(spec, project, session),
    "hellgraph:graph-query": lambda spec, project, session: _hellgraph_query(spec, project),
    "hellgraph:graph-stats": lambda spec, project, session: _hellgraph_stats(spec, project),
    "spark-runner": lambda spec, project, session: _spark(spec, project, session),
    "model-server": lambda spec, project, session: _inference(spec, project, session),
    "gateway:ingest": lambda spec, project, session: _ingest(spec, project, session),
    "gateway:parse": lambda spec, project, session: _parse(spec, project, session),
    "holmes": lambda spec, project, session: _extraction(spec, project, session),
    "open-data": lambda spec, project, session: _reconcile(spec, project, session),
    "sql": lambda spec, project, session: _sql_load(spec, project, session),
}


def set_backend(key: str, fn: Callable[..., Awaitable[AdapterResult]]) -> None:
    _BACKENDS[key] = fn


async def dispatch(kind: str, backend: str, spec: dict, project: str, session: str | None) -> AdapterResult:
    fn = _BACKENDS.get(f"{backend}:{kind}") or _BACKENDS.get(backend)
    if fn is None:
        return {"outputs": [], "runtime": backend, "status": "degraded", "error": None,
                "degraded": f"backend '{backend}' for kind '{kind}' not wired yet"}
    return await fn(spec, project, session)


# ── provenance write-back: every run becomes nodes+edges in the graph ──
async def write_provenance(delta: GraphDelta) -> bool:
    """Best-effort: persist the run's subgraph to hellgraph. Never fails the compute."""
    if not delta.nodes:
        return False
    try:
        async with httpx.AsyncClient(timeout=15.0) as c:
            for n in delta.nodes:
                await c.post(f"{HELLGRAPH_URL}/api/graph/node",
                             json={"id": n.id, "labels": n.labels, "properties": n.properties})
            for e in delta.edges:
                await c.post(f"{HELLGRAPH_URL}/api/graph/edge",
                             json={"label": e.label, "from": e.from_, "to": e.to, "properties": e.properties})
        return True
    except Exception:  # noqa: BLE001 — provenance write is best-effort
        return False


def build_delta(project: str, kind: str, backend: str, receipt_id: str, epistemic: str,
                inputs_sha: str | None = None, outputs_sha: str | None = None) -> GraphDelta:
    """The run's provenance subgraph, dual-labelled: OUR native labels (ComputeRun,
    Receipt, …) AND W3C PROV-O terms so it federates with any PROV-aware store —
    the run is a `prov:Activity`, the receipt/output/input are `prov:Entity`, and
    the edges are `prov:wasGeneratedBy` / `prov:used` / `prov:wasDerivedFrom`.
    """
    short = receipt_id.replace("sha256:", "")[:12]
    run_id = f"proj-{project}:compute:{short}"
    rc_id = f"proj-{project}:receipt:{short}"
    out_id = f"proj-{project}:output:{short}"
    in_id = f"proj-{project}:input:{short}"
    nodes = [
        GraphNode(id=run_id, labels=[project, "ComputeRun", kind, "prov:Activity"],
                  properties={"kind": kind, "backend": backend, "epistemic_mode": epistemic}),
        GraphNode(id=rc_id, labels=[project, "Receipt", "prov:Entity"],
                  properties={"receipt": receipt_id}),
        GraphNode(id=out_id, labels=[project, "ComputeOutput", "prov:Entity"],
                  properties={"outputs_sha": outputs_sha, "epistemic_mode": epistemic}),
        GraphNode(id=in_id, labels=[project, "ComputeInput", "prov:Entity"],
                  properties={"inputs_sha": inputs_sha}),
    ]
    edges = [
        # native label kept AND its PROV-O counterpart, side by side
        GraphEdge.model_validate({"label": "HAS_RECEIPT", "from": run_id, "to": rc_id}),
        GraphEdge.model_validate({"label": "prov:wasGeneratedBy", "from": rc_id, "to": run_id}),
        GraphEdge.model_validate({"label": "prov:wasGeneratedBy", "from": out_id, "to": run_id}),
        GraphEdge.model_validate({"label": "prov:used", "from": run_id, "to": in_id}),
        GraphEdge.model_validate({"label": "prov:wasDerivedFrom", "from": out_id, "to": in_id}),
    ]
    return GraphDelta(nodes=nodes, edges=edges)
