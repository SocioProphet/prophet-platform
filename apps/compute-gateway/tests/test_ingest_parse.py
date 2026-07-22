"""IFM stages 01–02: ingest (content-address the pack) + parse (bytes → page-keeping blocks),
completing the 5-stage governed pipeline: ingest → parse → extract → reconcile → load.

Documents are built in-test (python-pptx / pypdf), not fixtures — a parser bump that changes
behavior fails here first.
"""
import base64
import hashlib
import importlib
import io
import os

os.environ["GATEWAY_TOKEN"] = "t"
os.environ["COMPUTE_ENTITLEMENTS"] = "demo"
os.environ["GATEWAY_WRITE_PROVENANCE"] = "false"

from fastapi.testclient import TestClient  # noqa: E402

from compute_gateway import adapters, engine, receipts, server, zerotrust  # noqa: E402

importlib.reload(zerotrust)
importlib.reload(engine)
importlib.reload(server)
client = TestClient(server.app)
AUTH = {"Authorization": "Bearer t"}


_ORIG_HOLMES = adapters._BACKENDS["holmes"]
_ORIG_REFERENCE = adapters._REFERENCE


def setup_function():
    os.environ["COMPUTE_ENTITLEMENTS"] = "demo"
    receipts._CHAINS.clear()
    engine._MEMO.clear()


def teardown_function():
    # module-level injections must not leak into other test files (suite runs in one process)
    adapters.set_backend("holmes", _ORIG_HOLMES)
    adapters.set_reference_resolver(_ORIG_REFERENCE)


def _compute(body):
    return client.post("/v1/compute", json={"project": "demo", **body}, headers=AUTH).json()


def _b64(raw: bytes) -> str:
    return base64.b64encode(raw).decode()


def _pptx_pack() -> bytes:
    """A minimal graphics-heavy 'investor pack': title text + a financials table."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])   # title-only layout
    slide.shapes.title.text = "GYG FY26 Results"
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1)).table
    tbl.cell(0, 0).text = "Revenue"
    tbl.cell(0, 1).text = "$1,204m"
    tbl.cell(1, 0).text = "Net profit"
    tbl.cell(1, 1).text = "$14m"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── ingest ──
def test_ingest_content_addresses_the_raw_bytes():
    raw = b"Quarterly trading update.\n\nRevenue rose."
    r = _compute({"kind": "ingest", "spec": {"document_b64": _b64(raw), "filename": "update.txt"}})
    assert r["status"] == "ok" and r["epistemic_status"] == "observed"
    d = r["outputs"][0]["data"]
    assert d["sha256"] == hashlib.sha256(raw).hexdigest()      # hash of the BYTES, not the name
    assert d["size"] == len(raw) and d["media_type"] == "text/plain"
    assert r["receipt"]["id"]                                  # sealed like any governed compute


def test_ingest_rejects_missing_or_bad_b64():
    assert _compute({"kind": "ingest", "spec": {}})["status"] == "error"
    assert _compute({"kind": "ingest", "spec": {"document_b64": "!!not-b64!!"}})["status"] == "error"


def test_ingest_media_sniff_prefers_magic_bytes():
    pdfish = b"%PDF-1.7 minimal"
    r = _compute({"kind": "ingest", "spec": {"document_b64": _b64(pdfish), "filename": "lying.txt"}})
    assert r["outputs"][0]["data"]["media_type"] == "application/pdf"   # magic bytes beat the filename


# ── parse ──
def test_parse_pptx_keeps_pages_tables_and_bbox():
    r = _compute({"kind": "parse", "spec": {"document_b64": _b64(_pptx_pack()), "filename": "pack.pptx"}})
    assert r["status"] == "ok" and r["epistemic_status"] == "observed"
    d = r["outputs"][0]["data"]
    kinds = {b["kind"] for b in d["blocks"]}
    assert "table" in kinds and "text" in kinds
    table = next(b for b in d["blocks"] if b["kind"] == "table")
    assert "Revenue | $1,204m" in table["text"]                # tables carry the numbers
    assert table["page"] == 1 and table["bbox"] is not None    # page + region provenance
    assert all(isinstance(v, (int, float)) for v in table["bbox"])


def test_parse_text_and_error_paths():
    r = _compute({"kind": "parse", "spec": {"document_b64": _b64(b"Para one.\n\nPara two."), "filename": "n.txt"}})
    assert r["status"] == "ok" and len(r["outputs"][0]["data"]["blocks"]) == 2

    # blank PDF (scan-like) → honest error naming OCR, never a crash
    from pypdf import PdfWriter
    w = PdfWriter(); w.add_blank_page(width=612, height=792)
    buf = io.BytesIO(); w.write(buf)
    r = _compute({"kind": "parse", "spec": {"document_b64": _b64(buf.getvalue()), "filename": "scan.pdf"}})
    assert r["status"] == "error" and "OCR" in r["error"]

    # corrupt PDF → error, not a 500
    r = _compute({"kind": "parse", "spec": {"document_b64": _b64(b"%PDF-1.4 garbage"), "filename": "b.pdf"}})
    assert r["status"] == "error"


# ── the full 5-stage governed pipeline ──
def test_five_stage_pipeline_document_to_sql(tmp_path, monkeypatch):
    monkeypatch.setattr(adapters, "SQL_DSN", str(tmp_path / "ifm.db"), raising=False)

    async def ref(entity, field, period):
        return {"revenue": 1204.0}.get(field)                  # open-data stand-in agrees with the pack
    adapters.set_reference_resolver(ref)

    async def extract_from_blocks(spec, project, session):
        # a deterministic extractor over the THREADED parse blocks — proves data flows
        # ingest→parse→extract through `from`, not via hand-fed specs
        table_txt = next(b["text"] for b in spec["blocks"] if b["kind"] == "table")
        first = table_txt.splitlines()[0].split(" | ")
        value = float(first[1].replace("$", "").replace(",", "").rstrip("m"))
        return {"outputs": [adapters.ComputeOutput(type="table", data={
                    "table": spec["target_schema"]["table"],
                    "rows": [{"field": "revenue", "value": value, "unit": "AUD_m",
                              "page": spec["blocks"][0]["page"], "source_span": "tbl1/r1",
                              "warrant": "observed"}],
                    "entity": spec.get("entity"), "period": spec.get("period")})],
                "runtime": "holmes", "status": "ok", "error": None, "degraded": None,
                "epistemic": "observed"}
    adapters.set_backend("holmes", extract_from_blocks)

    r = _compute({"kind": "workflow", "spec": {"steps": [
        {"id": "ingest", "kind": "ingest", "spec": {"document_b64": _b64(_pptx_pack()), "filename": "pack.pptx"}},
        {"id": "parse", "kind": "parse", "from": "ingest"},
        {"id": "extract", "kind": "extraction", "from": "parse",
         "spec": {"target_schema": {"table": "financials"}, "entity": {"cik": "0"}, "period": "FY26"}},
        {"id": "reconcile", "kind": "reconcile", "from": "extract", "spec": {"tolerance": 0.01}},
        {"id": "load", "kind": "load", "from": "reconcile", "spec": {"table": "financials"}},
    ]}})
    assert r["status"] == "ok" and r["kind"] == "workflow"
    d = r["outputs"][0]["data"]
    steps = {s["id"]: s for s in d["steps"]}
    assert [*steps] == ["ingest", "parse", "extract", "reconcile", "load"]
    assert steps["ingest"]["epistemic_status"] == "observed"
    assert steps["reconcile"]["epistemic_status"] == "verified"   # pack agreed with the reference
    assert all(s["receipt"] for s in d["steps"])                  # every stage sealed its own receipt
    # 5 step receipts + 1 composite
    assert client.get("/v1/receipts", params={"project": "demo"}, headers=AUTH).json()["count"] == 6
